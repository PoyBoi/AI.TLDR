# [Basics]
import os, pathlib, hashlib, heapq, json, time, tempfile, subprocess, statistics
from typing import Optional, Literal
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor, as_completed
from multiprocessing import Manager
from tqdm import tqdm

# [Ingestion]
import PyPDF2, unstructured, orjson, io
from python_calamine import CalamineWorkbook
from charset_normalizer import from_path

# [LangChain]
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

# [LangGraph]

class modularRAG:
    def __init__(
            self,

            # RAG Core Args
            method: Literal["deep_agent", "basic_rag"] = "basic_rag", web_search: bool = True, chat_as_vdb: bool = True, re_ranking:bool = True, grading_step:bool = True,
            static_path:pathlib.Path = None, temp_path:pathlib.Path = None, # Assign temp path later as relative pathing

            # RAG Quality Args
            # [DEFAULTS]
            chunk_size:int = 400, chunk_overlap:int = 50, retrieval_method: Literal["sparse", "dense", "hybrid"] = "hybrid",
            text_splitter: Literal["document_based", "recursive_Character"] = "recursive_Character",

            # RAG QOL Args
            streaming: bool = True,
        ):
        # Initialising Args into Vars
        self.static_path = static_path
        self.temp_path = temp_path
        # Synchronus Tasks Below
        None
        # Multi-Processing Tasks Below
        # Will MP this later
        self.doc_loader()

    def doc_loader(self):
        # get doc's loc (can be temp path made by the code) -> load location of all the files -> batch divide (I need to find a good batch size, equally dividing it over 10 threads sounds fine) -|
        # |-> multiprocess/async(read based on file type + assign metadata (hash) -> convert to json / .md (need to look up the specifics) -|
        # |-> Make a tracker file that checks against the files to make sure a re-read isn't happening (compare hash as well as time modified against the one in the file)-> append to storage object)
        
        if self.static_path is not None:
            doc_loc = pathlib.Path(self.static_path)
        elif self.temp_path is not None:
            doc_loc = pathlib.Path(self.temp_path)
        else:
            EC = "Error - No path defined to work on"
            print(EC)
            return EC

        # Loads pre-existing cache file in process_cache.json
        cache_path = pathlib.Path("process_cache.json")
        if cache_path.exists():
            self.process_cache = json.loads(cache_path.read_text())
            # print(self.process_cache)
        else:
            self.process_cache = []

        # Build whole-file entries first — cache exclusion must run at file
        # granularity, before any PDF gets expanded into per-page entries,
        # or the same file gets cache-checked N times redundantly.
        whole_files = [
            {
                "file_path": f,
                "file_size": f.stat().st_size,
                "time_last_change": f.stat().st_ctime,
                "time_last_modification": f.stat().st_mtime,
                "file_hash": self._hash_file(f),
            } for f in doc_loc.iterdir() if f.is_file() and "placeholder" not in str(f)
        ]
        total_files = len(whole_files)

        if total_files == 0:
            print(f"No files found at {doc_loc}")
            return []

        # Exclusion + Updating Logic (multiprocessed, at whole-file granularity)
        cache_lookup = {str(item["file_path"]): item for item in (self.process_cache or [])}
        with ProcessPoolExecutor() as executor:
            files_to_process = list(
                executor.map(self.file_needs_update, whole_files, [cache_lookup] * total_files)
            )
        files_to_process = [f for f in files_to_process if f is not None]

        if not files_to_process:
            print("No files need processing — cache is up to date.")
            return []

        # Expands only the files that survived the cache check (with per-page entries for PDFs). Everything else passes through as-is.
        files_expanded = self._expand_pdfs_to_pages(files_to_process)

        batched_files = self.batcher(list_of_files=files_expanded)

        # Multiprocessing Manager for TQDM
        with Manager() as manager:
            progress_queue = manager.Queue()
            total_units = len(files_expanded)  # per-file/per-page granularity

            with ProcessPoolExecutor() as executor:
                futures = [
                    executor.submit(self.read_files, item, progress_queue)
                    for item in batched_files.items()
                ]

                with tqdm(total=total_units, desc="Reading files") as pbar:
                    while not all(f.done() for f in futures):
                        while not progress_queue.empty():
                            progress_queue.get()
                            pbar.update(1)
                        time.sleep(0.05)
                    # drain whatever landed after the last check
                    while not progress_queue.empty():
                        progress_queue.get()
                        pbar.update(1)

                processed_batches = [f.result() for f in futures]

        # Flatten {core_id, assignments} wrapper structure into a flat list of per-file/per-page result dicts before merging PDF pages back into single documents.
        flat_results = []
        for batch in processed_batches:
            for file_entry in batch["assignments"]:
                if file_entry["text_content"] is not None:
                    flat_results.append(file_entry["text_content"])

        merged_results = self.merge_pdf_page_results(flat_results)

        # print(merged_results)

        """
        Urgent TODO's

        - Append it to the VDB and cache it based on what time the VDB was made 
            - Check for parameters as well as best methods for storage
                - Make sure metadata and text splitting is working well
        - Add the agentic call methods first and then go with the linear flow
            - make sure to follow the 5-agent plan (supervisor -> guardrails -> RAG -> LLM -> Groundedness Checker -> Output)
                - guardrails - try:
                    - agentic routing LLM call
                    - small model that checks for malicious prompts
                    - check for off-domain queries
                - cross-check with the output validation loop and how it can be best done in a prod scenario
            - make sure the metadata is also returned
        - Add a swapper for AWS in the LLM section and make it configurable 
            - API key must exist inside of the `.env`
        - Add in-chat temp VDB
            - Check how to handle shortlong-term memory & how it's done in prod
                - Add the parent-child long-term memory format
                - Add fact summarisation for conversations above X tokens
                - Look at that new memory handling tool open-sourced by that actress
        - Visualise the vector database via an image / 3D space visualisation
            - Pre-render it so that it doesn't have to rendered in real time and save it
            - Make it viewable in a browser / pop up window spawned from the python process
        - Add ingestion stats from "analyze_chunks" to a .log / .txt file
            - Make a folder called "stats" (or something better idk) and name the file based on when the ingestion was carried out, if successful 
        - Find a way to monitor the input / output tokens and their attributed costs 
            - this needs to displayed in the front-end, need to be a part of the message metadata (will include previous and current token length and pricing based on the service used)
        - Implement a switch between chunking methods so it's easy for advanced users to change it and see which method will work the best
            - can also have a method that uses all the chunking method and shows which one can be the best one
                - MAYBE automatically use this method
        """

        content_to_write_to_json = self.process_cache + merged_results
        for r in content_to_write_to_json:
            # Chunking process applied to the data inside of results
            r["chunks"] = self.text_splitter(r["chunks"])
        # print(content_to_write_to_json)

        # Update the file with the new chunk data
        cache_path = pathlib.Path("test.json") #TODO: Remove this once testing is done
        try:
            with open(cache_path, "w+") as file:
                json.dump(content_to_write_to_json, file, indent=4)
        except Exception as e:
            print(f"Failed to write knowledge cache due to : {e}, re-trying...")
            try:
                with open("process_cache_backup.json", "w+") as file:
                    json.dump(content_to_write_to_json, file, indent=4)
            except Exception as e:
                print(f"Failed to write back-up knowledge cache as well due to {e}, skipping writing cache, reverting to older version...")
                with open(cache_path, "w+") as file:
                    json.dump(self.process_cache, file, indent=4)

        try:
            docs = [
                Document(page_content=r["text"], metadata={
                    "source": r["file_path"], 
                    "chunks": r["chunks"],
                    "file_size": r["file_size"],
                    "time_last_change": r["time_last_change"],
                    "time_last_modification": r["time_last_modification"],
                    "file_hash": r["file_hash"]
                    }
                )
                for r in content_to_write_to_json
            ]
            print("Documents updated to VDB-compliant structure...")
        except Exception as e:
            print(f"Failure of addition of documents to VDB-compliant structure, reason: {e}")

        # report = self.analyze_chunks(content_to_write_to_json)

        return merged_results

    def file_needs_update(self, file_data:dict = None, cache_lookup:dict = None):
        file_path = str(file_data["file_path"])
        cached_entry = (cache_lookup or {}).get(file_path)      # Finding the same entry as the one as which we have

        # Happens when the file exists in the cached file, aka it HAS been read before and is getting re-read
        if cached_entry is not None:
            print(f"File {file_path} already exists in cache, checking if it has been updated since...")

            # Comparing File Hashes first as that's a good way to check if the file's contents have been changed
            if file_data.get("file_hash"):
                if file_data["file_hash"] == cached_entry["file_hash"]:
                    print(f"File hash for file {file_path} is similar and doesn't need an update")
                    return None     # aka file doesn't need re-reading
                else:
                    return file_data        # Returning file to be re-read because the hash is dissimilar aka it has been changed since last read
            
            # Backup delta comparision to make sure the file hasn't been updated, if it has, add it to the list of files to be updated / re-read
            elif file_data["file_size"] == cached_entry["file_size"] or file_data["time_last_change"] == cached_entry["time_last_change"] or file_data["time_last_modification"] == cached_entry["time_last_modification"]:
                print(f"File Modification Data for file {file_path} is similar and doesn't need an update")
                return None         # aka file doesn't need re-reading
            else:
                return file_data        # Returning file to be re-read because the back-up data is dissimilar aka it has been changed since last read

        else:
            # Checks if the file itself even is valid
            if os.path.isfile(file_path):
                return file_data        # Returned where there is no cache'd entry, aka this is the first time the file is being read
            else:
                print(f"File {file_path} does not exist at it's path, please re-check, this file will not be read.")
                return None     # Returned when the file's path is incorrect / corrupted

    def batcher(self, list_of_files: list = None):
        # Creates batches based on file size and divides it based on the total core count of the user's CPU along with multiprocess pooling
        usable_cpu_count = max(1, (os.cpu_count() or 1) - 2)
        print(f"\nUsable CPU count for this session is: {usable_cpu_count}\n")

        # load list of dicts and and separate them into batches (size of file accumulated over X total cores)
        files_sorted = sorted(list_of_files, key=lambda f: f["file_size"], reverse=True)
    
        heap = [(0, i) for i in range(usable_cpu_count)]
        heapq.heapify(heap)
        
        assignments = {i: [] for i in range(usable_cpu_count)}
        
        for i in files_sorted:
            size = i["file_size"]
            load, core_id = heapq.heappop(heap)
            assignments[core_id].append(i)
            heapq.heappush(heap, (load + size, core_id))
        
        return assignments

    def read_files(self, file_data: tuple = None, progress_queue=None) -> dict:
        core_id, assignments = file_data

        for file in assignments:
            output = self.read_file(
                file_loc=file["file_path"],
                page_index=file.get("page_index"),
                num_pages=file.get("num_pages"),
                file_size=file.get("file_size"),
                time_last_change=file.get("time_last_change"),
                time_last_modification=file.get("time_last_modification"),
                file_hash=file.get("file_hash"),
            )
            file["text_content"] = output
            if progress_queue is not None: # Adding descrete counter for TQDM's bar
                progress_queue.put(1)

        return {"core_id": core_id, "assignments": assignments}

    def read_file(
            self, file_loc:pathlib.Path = None, page_index: int = None, num_pages: int = None,
            file_size: int = None, time_last_change: float = None, time_last_modification: float = None, file_hash: str = None
        ) -> None:

        action_map = {
            ".pdf": self.read_pdf,
            ".docx": self.read_doc,
            ".doc": self.deprecated_file_format,
            ".xlsx": self.read_excel,
            ".xls": self.read_excel,
            ".md": self.read_md,
            ".json": self.read_json,
            ".png": self.read_img,
            ".jpg": self.read_img,
            ".jpeg": self.read_img,
            ".webp": self.read_img,
            ".py": self.read_code,
            ".js": self.read_code,
            ".html": self.read_code,
            ".css": self.read_code,
            ".ts": self.read_code,
            ".sh": self.read_code,
            ".bat": self.read_code,
            ".txt": self.read_txt,
            ".ppt": self.deprecated_file_format,
            ".pptx": self.read_ppt,
        }
        # Does the actual reading, done for the sake of modularisation and code-upkeep
        # print(file_loc)
        file_format = "".join(file_loc.suffixes)
        # print(file_format)
        file_read_function = action_map.get(file_format, self.unsupported_file_format)

        file_meta = {
            "file_size": file_size,
            "time_last_change": time_last_change,
            "time_last_modification": time_last_modification,
            "file_hash": file_hash,
        }

        try:
            if file_format == ".pdf":
                result = self.read_pdf(
                    file_loc=file_loc,
                    file_format=file_format,
                    page_index=page_index,
                    num_pages=num_pages,
                )
            else:
                result = file_read_function(file_loc=file_loc, file_format=file_format)

            if isinstance(result, dict):
                result.update(file_meta)
            return result

        except Exception as e:
            print(f"Failed to read {file_loc} ({file_format}, page_index={page_index}): {e}")
            return {
                "file_path": str(file_loc),
                "file_format": str(file_format),
                "page_index": page_index,
                "text": "",
                "chunks": [],
                "error": str(e),
                **file_meta,
            }

    def read_doc(self, file_loc: pathlib.Path = None, file_format: str = None):
        from docx import Document as DocxDocument
        print(f"Reading DOC/DOCX with file format {file_format}")

        if file_format != ".docx":
            raise ValueError(f"Unsupported extension for read_doc: {file_format}")

        doc = DocxDocument(file_loc)
        chunks = []
        full_text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                chunks.append({
                    "type": "paragraph",
                    "text": para.text,
                    "section_path": [],
                    "page": None,
                })
                full_text_parts.append(para.text)

        for t_idx, table in enumerate(doc.tables, start=1):
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            markdown = self._table_to_markdown(rows)
            chunks.append({
                "type": "table",
                "text": markdown,
                "rows": rows,
                "section_path": [f"Table {t_idx}"],
                "page": None,
            })
            full_text_parts.append(markdown)

        return {
            "file_path": str(file_loc),
            "file_format": file_format,
            "text": "\n".join(full_text_parts),
            "chunks": chunks,
        }

    def read_excel(self, file_loc: pathlib.Path = None, file_format: str = None):
        print(f"Reading XLSX/XLS with file format {file_format}")
        workbook = CalamineWorkbook.from_path(str(file_loc))

        chunks = []
        full_text_parts = []

        for sheet_name in workbook.sheet_names:
            rows = workbook.get_sheet_by_name(sheet_name).to_python()
            rows = [[str(cell) if cell is not None else "" for cell in row] for row in rows]
            markdown = self._table_to_markdown(rows)
            chunks.append({
                "type": "table",
                "text": markdown,
                "rows": rows,
                "section_path": [f"Sheet: {sheet_name}"],
                "page": None,
            })
            full_text_parts.append(f"# Sheet: {sheet_name}\n{markdown}")

        return {
            "file_path": str(file_loc),
            "file_format": file_format,
            "text": "\n\n".join(full_text_parts),
            "chunks": chunks,
        }
    
    def read_md(self, file_loc: pathlib.Path = None, file_format: str = None):
        print(f"Reading MarkDown with file format {file_format}")
        text = file_loc.read_text(encoding="utf-8", errors="replace")

        return {
            "file_path": str(file_loc),
            "file_format": file_format,
            "text": text,
            "chunks": [{
                "type": "paragraph",
                "text": text,
                "section_path": [],
                "page": None,
            }],
        }

    def read_json(self, file_loc: pathlib.Path = None, file_format: str = None):
        print(f"Reading JSON with file format {file_format}")
        with open(file_loc, "rb") as f:
            data = orjson.loads(f.read())
        text = orjson.dumps(data, option=orjson.OPT_INDENT_2).decode("utf-8")

        return {
            "file_path": str(file_loc),
            "file_format": file_format,
            "text": text,
            "chunks": [{
                "type": "paragraph",
                "text": text,
                "section_path": [],
                "page": None,
            }],
        }

    def read_code(self, file_loc: pathlib.Path = None, file_format: str = None):
        print(f"Reading Code File with file format {file_format}")
        try:
            text = file_loc.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            result = from_path(str(file_loc)).best()
            if result is None:
                raise ValueError(f"Could not decode {file_loc}")
            text = str(result)

        return {
            "file_path": str(file_loc),
            "file_format": file_format,
            "text": text,
            "chunks": [{
                "type": "code",
                "text": text,
                "section_path": [],
                "page": None,
            }],
        }

    def get_pdf_page_count(self, file_loc: pathlib.Path) -> int:
        import pymupdf
        doc = pymupdf.open(str(file_loc))
        count = len(doc)
        doc.close()
        return count

    def _expand_pdfs_to_pages(self, files: list) -> list:
        expanded = []

        for entry in files:
            file_path = entry["file_path"]
            suffix = file_path.suffix.lower()

            if suffix != ".pdf":
                expanded.append({**entry, "page_index": None, "num_pages": None})
                continue

            num_pages = self.get_pdf_page_count(file_path)
            # Weight used by the batcher for load-balancing — dividing the
            # file's real size across its pages gives a rough per-page cost,
            # which is a better proxy than treating the whole file as one unit.
            per_page_weight = max(1, entry["file_size"] // max(1, num_pages))

            for page_index in range(num_pages):
                expanded.append({
                    **entry,
                    "file_size": per_page_weight,
                    "page_index": page_index,
                    "num_pages": num_pages,
                })

        return expanded

    def merge_pdf_page_results(self, all_results: list) -> list:
        from collections import defaultdict

        grouped = defaultdict(list)
        passthrough_results = []

        for r in all_results:
            if r.get("page_index") is not None:
                grouped[r["file_path"]].append(r)
            else:
                passthrough_results.append(r)

        merged = []
        for source_file, page_results in grouped.items():
            page_results.sort(key=lambda r: r["page_index"])
            merged.append({
                "file_path": source_file,
                "file_format": ".pdf",
                "text": "\n\n".join(r["text"] for r in page_results if r["text"]),
                "chunks": [c for r in page_results for c in r["chunks"]],
                "file_size": page_results[0].get("file_size"),
                "time_last_change": page_results[0].get("time_last_change"),
                "time_last_modification": page_results[0].get("time_last_modification"),
                "file_hash": page_results[0].get("file_hash"),
            })

        return passthrough_results + merged

    def read_pdf(self, file_loc: pathlib.Path = None, file_format: str = None, page_index: int = None, num_pages: int = None):
        # Need to optimise this code
        import pymupdf
        import pdfplumber

        # Commented this out for sake of new TQDM progress bar
        # if page_index is not None:
        #     print(f"Reading PDF page {page_index + 1} of {file_loc}")
        # else:
        #     print(f"Reading PDF (whole file) {file_loc}")

        doc = pymupdf.open(str(file_loc))
        target_indices = [page_index] if page_index is not None else list(range(len(doc)))

        # Pull tables only for the target page(s) via pdfplumber.
        tables_by_page = {}
        with pdfplumber.open(str(file_loc)) as plumber_pdf:
            for idx in target_indices:
                if idx >= len(plumber_pdf.pages):
                    continue
                extracted = plumber_pdf.pages[idx].extract_tables()
                if extracted:
                    tables_by_page[idx + 1] = extracted

        chunks = []
        full_text_parts = []

        for idx in target_indices:
            page = doc[idx]
            page_num = idx + 1

            page_text = page.get_text("text").strip()

            if not page_text:
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("png")
                ocr_text = self.read_img(
                    img_bytes=img_bytes,
                    source_hint=f"{file_loc.stem}_p{page_num}_scan",
                )
                chunks.append({
                    "type": "paragraph",
                    "text": ocr_text,
                    "section_path": [f"Page {page_num}"],
                    "page": page_num,
                    "source": "ocr",
                })
                full_text_parts.append(ocr_text)
            else:
                chunks.append({
                    "type": "paragraph",
                    "text": page_text,
                    "section_path": [f"Page {page_num}"],
                    "page": page_num,
                    "source": "text_layer",
                })
                full_text_parts.append(page_text)

            for table_idx, table_rows in enumerate(tables_by_page.get(page_num, []), start=1):
                clean_rows = [
                    [cell if cell is not None else "" for cell in row]
                    for row in table_rows
                ]
                markdown = self._table_to_markdown(clean_rows)
                chunks.append({
                    "type": "table",
                    "text": markdown,
                    "rows": clean_rows,
                    "section_path": [f"Page {page_num}", f"Table {table_idx}"],
                    "page": page_num,
                })
                full_text_parts.append(markdown)

            for img_idx, img in enumerate(page.get_images(full=True), start=1):
                xref = img[0]
                try:
                    pix = pymupdf.Pixmap(doc, xref)
                    if pix.n - pix.alpha > 3:
                        pix = pymupdf.Pixmap(pymupdf.csRGB, pix)
                    img_bytes = pix.tobytes("png")
                except Exception as e:
                    print(f"Skipping image xref {xref} on page {page_num}: {e}")
                    continue

                img_desc = self.read_img(
                    img_bytes=img_bytes,
                    source_hint=f"{file_loc.stem}_p{page_num}_img{img_idx}",
                )
                chunks.append({
                    "type": "image_ref",
                    "text": img_desc,
                    "section_path": [f"Page {page_num}", f"Image {img_idx}"],
                    "page": page_num,
                    "xref": xref,
                })

        doc.close()

        return {
            "file_path": str(file_loc),
            "file_format": file_format,
            "page_index": page_index,   # None means this result is the whole file
            "text": "\n\n".join(full_text_parts),
            "chunks": chunks,
        }


    def read_img(
        self, file_loc: pathlib.Path = None, file_format: str = None, img_bytes: bytes = None, source_hint: str = None,
    ):
        """
        Two call sites:
        - Standalone image file: file_loc + file_format set, img_bytes is None.
        - Embedded image from read_pdf: img_bytes + source_hint set, file_loc is None.
        """
        from PIL import Image
        import pytesseract

        if img_bytes is not None:
            img = Image.open(io.BytesIO(img_bytes))
            stem = source_hint or "image"
            source_ref = f"embedded:{stem}"
        else:
            print(f"Reading Image with file format {file_format}")
            img = Image.open(file_loc)
            stem = file_loc.stem
            source_ref = str(file_loc)

        out_dir = pathlib.Path("extracted_images")
        out_dir.mkdir(exist_ok=True)
        out_path = out_dir / f"{stem}.png"
        img.save(out_path)

        try:
            ocr_text = pytesseract.image_to_string(img).strip()
        except Exception as e:
            print(f"OCR failed for {stem}: {e}")
            ocr_text = ""

        text = f"[image: {out_path}]\n{ocr_text}" if ocr_text else f"[image: {out_path}] (no extractable text — caption pending)"

        # Only wrap in the full dict shape when called as a top-level
        # dispatch reader (file_loc set). When called internally from
        # read_pdf, just return the description string as before.
        if file_loc is None:
            return text

        return {
            "file_path": source_ref,
            "file_format": file_format,
            "text": text,
            "chunks": [{
                "type": "image_ref",
                "text": text,
                "section_path": [],
                "page": None,
            }],
        }

    def _hash_file(self, path: pathlib.Path, chunk_size: int = 1 << 20) -> str:
        hasher = hashlib.sha256()
        with open(path, "rb") as fp:
            for block in iter(lambda: fp.read(chunk_size), b""):
                hasher.update(block)
        return hasher.hexdigest()

    def _table_to_markdown(self, rows: list) -> str:
        if not rows:
            return ""
        header, *body = rows
        md = ["| " + " | ".join(str(c) for c in header) + " |"]
        md.append("|" + "---|" * len(header))
        for row in body:
            md.append("| " + " | ".join(str(c) for c in row) + " |")
        return "\n".join(md)

    def read_txt(self, file_loc: pathlib.Path = None, file_format: str = None):
        print(f"Reading Text File with file format {file_format}")
        try:
            text = file_loc.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            result = from_path(str(file_loc)).best()
            if result is None:
                raise ValueError(f"Could not decode {file_loc}")
            text = str(result)

        return {
            "file_path": str(file_loc),
            "file_format": file_format,
            "text": text,
            "chunks": [{
                "type": "paragraph",
                "text": text,
                "section_path": [],
                "page": None,
            }],
        }

    def read_ppt(self, file_loc: pathlib.Path = None, file_format: str = None):
        print(f"Reading PPT/PPTX with file format {file_format}")
        if file_format != ".pptx":
            raise ValueError(f"Unsupported extension for read_ppt: {file_format}")

        from pptx import Presentation
        prs = Presentation(file_loc)

        chunks = []
        full_text_parts = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text.strip():
                            slide_parts.append(text)
                if shape.has_table:
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    markdown = self._table_to_markdown(rows)
                    chunks.append({
                        "type": "table",
                        "text": markdown,
                        "rows": rows,
                        "section_path": [f"Slide {slide_idx}"],
                        "page": slide_idx,
                    })
                    slide_parts.append(markdown)

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                note_text = slide.notes_slide.notes_text_frame.text
                slide_parts.append(f"[notes] {note_text}")

            slide_text = "\n".join(slide_parts)
            chunks.append({
                "type": "paragraph",
                "text": slide_text,
                "section_path": [f"Slide {slide_idx}"],
                "page": slide_idx,
            })
            full_text_parts.append(f"# Slide {slide_idx}\n{slide_text}")

        return {
            "file_path": str(file_loc),
            "file_format": file_format,
            "text": "\n\n".join(full_text_parts),
            "chunks": chunks,
        }

    def unsupported_file_format(self, file_loc: pathlib.Path = None, file_format: str = None) -> None:
        print(f"Provided file at {file_loc} does not have a valid file format")
        return None

    def deprecated_file_format(self, file_loc: pathlib.Path = None, file_format:str = None) -> None:
        print(f"Given file format {file_format} is deprecated and is not supported")

    def _build_breadcrumb(self, file_loc: pathlib.Path, section_path: list) -> str:
        doc_name = file_loc.stem if isinstance(file_loc, pathlib.Path) else pathlib.Path(str(file_loc)).stem
        parts = [doc_name] + [p for p in (section_path or []) if p]
        return "[" + " > ".join(parts) + "]"

    def text_splitter(self, chunks: list, chunk_size: int = 400, chunk_overlap: int = 50) -> list:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""],
        )

        result = []
        for c in chunks:
            if c["type"] != "paragraph" or len(c["text"]) <= chunk_size:
                result.append(c)
                continue

            pieces = splitter.split_text(c["text"])
            for i, piece in enumerate(pieces):
                result.append({
                    **c,                       # keeps type, section_path, page, source
                    "text": piece,
                    "split_index": i,          # marks it as a sub-chunk of the original
                    "split_of": len(pieces),
                })

        return result

    def analyze_chunks(self, content_to_write_to_json: list) -> dict:
        char_lens = []
        word_lens = []
        by_type = {}          # type -> list of char lengths
        by_source = {}         # file_path -> list of char lengths
        empty_chunks = 0
        tiny_chunks = []       # < 50 chars — likely noise / over-fragmented
        huge_chunks = []       # > 2000 chars — likely under-fragmented for embedding

        for r in content_to_write_to_json:
            source = r.get("file_path", "unknown")
            for c in r.get("chunks", []):
                text = (c.get("text") or "").strip()
                c_type = c.get("type", "unknown")
                n_chars = len(text)
                n_words = len(text.split())

                if n_chars == 0:
                    empty_chunks += 1
                    continue

                char_lens.append(n_chars)
                word_lens.append(n_words)
                by_type.setdefault(c_type, []).append(n_chars)
                by_source.setdefault(source, []).append(n_chars)

                if n_chars < 50:
                    tiny_chunks.append((source, c_type, n_chars))
                elif n_chars > 2000:
                    huge_chunks.append((source, c_type, n_chars))

        if not char_lens:
            print("No non-empty chunks found.")
            return {}

        report = {
            "total_chunks": len(char_lens),
            "empty_chunks": empty_chunks,
            "char_mean": statistics.mean(char_lens),
            "char_median": statistics.median(char_lens),
            "char_stdev": statistics.stdev(char_lens) if len(char_lens) > 1 else 0,
            "char_min": min(char_lens),
            "char_max": max(char_lens),
            "word_mean": statistics.mean(word_lens),
            "word_median": statistics.median(word_lens),
            "tiny_chunk_count": len(tiny_chunks),
            "huge_chunk_count": len(huge_chunks),
        }

        print("=== Chunk Size Report ===")
        print(f"Total chunks:        {report['total_chunks']}")
        print(f"Empty chunks:        {report['empty_chunks']}")
        print(f"Char length  - mean: {report['char_mean']:.1f}  median: {report['char_median']:.1f}  "
            f"stdev: {report['char_stdev']:.1f}  min: {report['char_min']}  max: {report['char_max']}")
        print(f"Word count   - mean: {report['word_mean']:.1f}  median: {report['word_median']:.1f}")
        print(f"Tiny chunks (<50 chars):  {report['tiny_chunk_count']}  ({report['tiny_chunk_count']/report['total_chunks']*100:.1f}%)")
        print(f"Huge chunks (>2000 chars): {report['huge_chunk_count']}  ({report['huge_chunk_count']/report['total_chunks']*100:.1f}%)")

        print("\n--- By chunk type ---")
        for c_type, lens in sorted(by_type.items(), key=lambda x: -len(x[1])):
            print(f"  {c_type:12s} count={len(lens):5d}  mean={statistics.mean(lens):7.1f}  "
                f"median={statistics.median(lens):7.1f}  max={max(lens):7d}")

        print("\n--- Worst offenders (files with the most extreme avg chunk size) ---")
        source_avgs = sorted(
            ((src, statistics.mean(lens), len(lens)) for src, lens in by_source.items()),
            key=lambda x: x[1],
        )
        print("  Smallest avg chunks:")
        for src, avg, n in source_avgs[:5]:
            print(f"    {src}  avg={avg:.1f}  n_chunks={n}")
        print("  Largest avg chunks:")
        for src, avg, n in source_avgs[-5:]:
            print(f"    {src}  avg={avg:.1f}  n_chunks={n}")

        return report

#__main__
if __name__ == "__main__":
    #TODO: This depends on where the file is being executed from, will need to change in the final edit
    RAG = modularRAG(static_path=r"/home/poyboi/VSCodesWSL/projects/AI.TLDR/input_folder")
    # RAG = modularRAG(static_path=r"../input_folder")